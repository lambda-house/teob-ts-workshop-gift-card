#!/usr/bin/env python3
"""
Build the complete DDD Workshop PPTX deck.
Run: /tmp/pptx-env/bin/python3 slides/build-slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme constants ───────────────────────────────────────────────────────

BG_LIGHT = RGBColor(0xEC, 0xEC, 0xEC)
BG_DARK = RGBColor(0x1E, 0x1E, 0x2E)

COLOR_DARK = RGBColor(0x22, 0x22, 0x22)
COLOR_BLUE = RGBColor(0x22, 0x66, 0xAA)
COLOR_GREEN = RGBColor(0x66, 0xBB, 0x6A)
COLOR_RED = RGBColor(0xEF, 0x53, 0x50)
COLOR_GRAY = RGBColor(0x44, 0x44, 0x44)
COLOR_LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ORANGE = RGBColor(0xFF, 0xA7, 0x26)
COLOR_PURPLE = RGBColor(0x9C, 0x27, 0xB0)

SLIDE_W = 12191695
SLIDE_H = 6858000

M_LEFT = Emu(731520)
M_TOP = Emu(548640)
M_RIGHT = Emu(731520)
CONTENT_W = SLIDE_W - M_LEFT - M_RIGHT

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK_LAYOUT = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, bold=False,
             color=COLOR_DARK, font_name=None, alignment=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    if font_name:
        p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_title(slide, text, color=COLOR_DARK, font_size=36):
    add_text(slide, M_LEFT, M_TOP, CONTENT_W, Emu(800000),
             text, font_size=font_size, bold=True, color=color)

def add_bullet_block(slide, items, top, color=COLOR_DARK, font_size=18,
                     bold=False, spacing=Emu(380000), left=None, width=None):
    if left is None:
        left = M_LEFT
    if width is None:
        width = CONTENT_W
    for i, item in enumerate(items):
        y = top + i * spacing
        prefix = "• " if not item.startswith(("─", "→", " ", "•")) else ""
        add_text(slide, left, y, width, spacing,
                 f"{prefix}{item}", font_size=font_size, bold=bold, color=color)

def add_code_box(slide, text, left, top, width, height, font_size=14):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    shape.line.fill.background()
    shape.adjustments[0] = 0.02
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(16)
    tf.margin_top = Pt(12)
    tf.margin_right = Pt(16)
    tf.margin_bottom = Pt(12)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "SF Mono"
    return shape

def add_accent_line(slide, top, width=None, color=COLOR_BLUE):
    if width is None:
        width = CONTENT_W
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, M_LEFT, top, width, Emu(36000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_box(slide, left, top, width, height, fill_color, text="",
            font_size=16, text_color=COLOR_WHITE, bold=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_arrow(slide, left, top, width, height, color=COLOR_GRAY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_accent_line(slide, Emu(2600000))
add_text(slide, M_LEFT, Emu(2800000), CONTENT_W, Emu(1200000),
         "DDD WITH TEOB", font_size=54, bold=True, color=COLOR_DARK)
add_text(slide, M_LEFT, Emu(4000000), CONTENT_W, Emu(600000),
         "Hands-On Workshop: Event-Sourced Aggregates in TypeScript",
         font_size=22, color=COLOR_GRAY)
add_text(slide, M_LEFT, Emu(5200000), CONTENT_W, Emu(400000),
         "Tim Evdokimov", font_size=18, color=COLOR_LIGHT_GRAY)
add_notes(slide,
    "TIMING: 0:00–0:05\n"
    "Welcome everyone. Before slides — ask the room:\n"
    "\"What broke in your last distributed system?\"\n"
    "Let 2-3 people share war stories. This primes them for why DDD + ES matters.\n"
    "Then: \"Today you'll build a complete event-sourced gift card service.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2: What You'll Build
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "WHAT YOU'LL BUILD TODAY")
add_accent_line(slide, Emu(1200000))

items = [
    ("1", "A DDD aggregate", "commands, events, state, invariants — pure functions", COLOR_BLUE),
    ("2", "An HTTP service", "your aggregate becomes a live API in one line", COLOR_BLUE),
    ("3", "A read model", "project events into query-optimized views", COLOR_GREEN),
    ("4", "LLM integration", "ctx.sync() calls OpenRouter — same pattern as any external API", COLOR_PURPLE),
]
for i, (num, title, subtitle, color) in enumerate(items):
    y = Emu(1600000 + i * 1250000)
    add_box(slide, M_LEFT, y, Emu(700000), Emu(700000), color,
            text=num, font_size=32, text_color=COLOR_WHITE)
    add_text(slide, Emu(1700000), y + Emu(50000), Emu(9000000), Emu(400000),
             title, font_size=24, bold=True, color=COLOR_DARK)
    add_text(slide, Emu(1700000), y + Emu(400000), Emu(9000000), Emu(300000),
             subtitle, font_size=16, color=COLOR_GRAY)

add_notes(slide,
    "TIMING: 0:05–0:08\n"
    "\"One domain — gift cards. Four layers of capability.\"\n"
    "\"You'll write code for 1 and 3. I'll demo 2 and 4.\"\n"
    "\"By the end you'll have a running HTTP service with AI integration.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3: Audience Calibration
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "QUICK CHECK — SHOW OF HANDS")
add_accent_line(slide, Emu(1200000))

questions = [
    "Who has implemented a DDD aggregate before?",
    "Who has worked with event sourcing in production?",
    "Who writes TypeScript regularly?",
    "Who has used Akka, Axon, or EventStoreDB?",
]
for i, q in enumerate(questions):
    y = Emu(1700000 + i * 900000)
    add_text(slide, Emu(1000000), y, Emu(10000000), Emu(600000),
             f"✋  {q}", font_size=22, color=COLOR_DARK)

add_notes(slide,
    "TIMING: 0:08–0:10\n"
    "This calibrates your depth.\n"
    "If most have DDD experience → skip basics, go deeper on TEOB specifics.\n"
    "If Akka/Axon users → draw parallels throughout."
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4: Road Map
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "THE ROAD MAP")
add_accent_line(slide, Emu(1200000))

stages = [
    ("Theory:\nES + DDD", COLOR_LIGHT_GRAY, "20 min"),
    ("Event\nStorming", COLOR_ORANGE, "10 min"),
    ("Exercise 1:\nAggregate", COLOR_GREEN, "30 min"),
    ("Demo:\nHTTP + API", COLOR_BLUE, "5 min"),
    ("Exercise 2:\nProjection", COLOR_GREEN, "20 min"),
    ("Exercise 3:\nAI Frontend", COLOR_GREEN, "10 min"),
    ("Demo:\nLLM", COLOR_PURPLE, "10 min"),
    ("Next\nSteps", COLOR_ORANGE, "5 min"),
]
box_w = Emu(1200000)
gap = Emu(80000)
total_w = len(stages) * box_w + (len(stages) - 1) * gap
start_x = Emu((SLIDE_W - total_w) // 2)
for i, (label, color, time) in enumerate(stages):
    x = start_x + i * (box_w + gap)
    add_box(slide, x, Emu(2200000), box_w, Emu(1800000), color,
            text=label, font_size=11, text_color=COLOR_WHITE)
    add_text(slide, x, Emu(4200000), box_w, Emu(400000),
             time, font_size=11, color=COLOR_LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

add_text(slide, M_LEFT, Emu(5200000), CONTENT_W, Emu(600000),
         "~2.5 hours  •  One domain (gift cards)  •  Three exercises + two live demos",
         font_size=16, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "TIMING: 0:10–0:12\n"
    "\"Green blocks = you write code. Blue/purple = I demo live.\"\n"
    "\"One domain all the way through — each layer builds on the last.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5: Traditional Architecture
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "TRADITIONAL ARCHITECTURE")
add_accent_line(slide, Emu(1200000))

add_text(slide, M_LEFT, Emu(1600000), Emu(5500000), Emu(400000),
         "Primary state: in database", font_size=22, bold=True, color=COLOR_DARK)

items = [
    "Database as source of truth: updates in place",
    "Many clients, larger states, frequent updates",
    "Locks and I/O overhead",
    "Single master doesn't scale",
    "Caching??...",
]
add_bullet_block(slide, items, Emu(2200000), color=COLOR_GRAY, font_size=17)

add_notes(slide,
    "TIMING: 0:12–0:16\n"
    "\"This is what most of us build day to day.\"\n"
    "Don't dwell — this is the setup for why ES is different."
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6: What Makes an Ideal Backend?
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "WHAT MAKES AN IDEAL BACKEND?")
add_accent_line(slide, Emu(1200000))

qualities = [
    ("(Almost) infinitely scalable", "No single master, no locks needed"),
    ("Lowest latency possible", "Primary state in memory, not on disk"),
    ("Always deterministic", "Durability, ordering, integrity guarantees"),
    ("Transparent, auditable logics", "Every state change is an event you can replay"),
    ("Below complexity barrier", "Express intent, not infrastructure"),
]
for i, (title, desc) in enumerate(qualities):
    y = Emu(1700000 + i * 850000)
    add_text(slide, Emu(1000000), y, Emu(5500000), Emu(400000),
             title, font_size=21, bold=True, color=COLOR_BLUE)
    add_text(slide, Emu(6800000), y + Emu(30000), Emu(5000000), Emu(400000),
             desc, font_size=16, color=COLOR_GRAY)

add_notes(slide, "TIMING: 0:16–0:18\n\"What if we could have ALL of these?\"")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7: Complete Behaviour
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "COMPLETE BEHAVIOUR: DATA + LOGICS")
add_accent_line(slide, Emu(1200000))

add_text(slide, M_LEFT, Emu(1600000), CONTENT_W, Emu(500000),
         "Business entities: orders, transactions, customers, digital twins...",
         font_size=19, color=COLOR_GRAY)
add_text(slide, M_LEFT, Emu(2300000), CONTENT_W, Emu(400000),
         "Own state and interactions with the world — fully covered:",
         font_size=19, bold=True, color=COLOR_DARK)

items = [
    "Incoming messages (commands) + immediate replies",
    "State changes (events applied to state)",
    "External effect calls (to other systems, LLMs, databases...)",
    "Deferred replies and scheduled timers",
]
add_bullet_block(slide, items, Emu(3000000), color=COLOR_DARK, font_size=18,
                 spacing=Emu(500000))

add_notes(slide, "TIMING: 0:18–0:20\n\"You'll implement exactly this pattern in 15 minutes.\"")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8: Event Sourcing Concept
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "EVENT SOURCING: THE IDEA")
add_accent_line(slide, Emu(1200000))

add_text(slide, M_LEFT, Emu(1600000), Emu(5000000), Emu(400000),
         "Primary state: in memory", font_size=22, bold=True, color=COLOR_DARK)

left_items = [
    "Log of updates (journal) as source of truth",
    "Stream — append-only, immutable",
    "Never in-place updates",
    "No locks, no I/O overhead — scales!",
]
add_bullet_block(slide, left_items, Emu(2200000), color=COLOR_GRAY, font_size=17)

journal_x = Emu(7500000)
add_text(slide, journal_x, Emu(1500000), Emu(4000000), Emu(400000),
         "Event Journal", font_size=16, bold=True, color=COLOR_BLUE,
         alignment=PP_ALIGN.CENTER)

events = ["CardIssued", "CardRedeemed", "CardRedeemed", "CardCancelled"]
for i, evt in enumerate(events):
    y = Emu(2000000 + i * 550000)
    add_box(slide, journal_x + Emu(200000), y, Emu(3200000), Emu(420000),
            RGBColor(0x2D, 0x2D, 0x2D), text=f"#{i+1}  {evt}", font_size=13,
            text_color=COLOR_GREEN)

add_text(slide, journal_x, Emu(4300000), Emu(3600000), Emu(400000),
         "→ append only, never mutate",
         font_size=14, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "TIMING: 0:20–0:23\n"
    "\"Instead of updating a row, we append facts.\"\n"
    "Draw parallel to git: commits are events, working directory is state."
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9: Pure Event Sourcing — decide/apply
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "PURE EVENT SOURCING: THE CORE LOOP")
add_accent_line(slide, Emu(1200000))

add_code_box(slide,
    "decide(State, Command) → Effect<Event, Reply>\n\n"
    "  → validate command against current state\n"
    "  → produce events (or reject)\n"
    "  → optionally reply to caller\n\n"
    "apply(State, Event) → State\n\n"
    "  → pure fold: no logic, just data transformation\n"
    "  → used for replay",
    M_LEFT, Emu(1600000), Emu(6000000), Emu(3400000), font_size=17)

flow_x = Emu(7500000)
add_box(slide, flow_x, Emu(1600000), Emu(2400000), Emu(600000),
        COLOR_BLUE, "Command", font_size=18)
add_text(slide, flow_x + Emu(900000), Emu(2300000), Emu(600000), Emu(400000),
         "↓", font_size=28, bold=True, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)
add_box(slide, flow_x, Emu(2700000), Emu(2400000), Emu(600000),
        COLOR_DARK, "decide()", font_size=18, text_color=COLOR_WHITE)
add_text(slide, flow_x + Emu(900000), Emu(3400000), Emu(600000), Emu(400000),
         "↓", font_size=28, bold=True, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)
add_box(slide, flow_x, Emu(3800000), Emu(2400000), Emu(600000),
        COLOR_GREEN, "Event(s)", font_size=18, text_color=COLOR_WHITE)
add_text(slide, flow_x + Emu(900000), Emu(4500000), Emu(600000), Emu(400000),
         "↓", font_size=28, bold=True, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)
add_box(slide, flow_x, Emu(4900000), Emu(2400000), Emu(600000),
        COLOR_ORANGE, "apply() → State'", font_size=18, text_color=COLOR_WHITE)

add_notes(slide,
    "TIMING: 0:23–0:27\n"
    "THIS IS THE CORE MENTAL MODEL. Spend time here.\n"
    "\"decide() is WHERE your business logic lives.\"\n"
    "\"apply() is PURE — no logic, just state transitions.\"\n"
    "Ask: \"Why can't we put logic in apply?\" → Because replay would re-validate."
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10: Effects
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "EVENT SOURCING: EFFECTS")
add_accent_line(slide, Emu(1200000))

add_text(slide, M_LEFT, Emu(1600000), CONTENT_W, Emu(400000),
         "decide() returns an Effect describing what should happen:",
         font_size=18, color=COLOR_GRAY)

effects = [
    ("persist(event)", "Store event(s) in the journal", COLOR_BLUE),
    ("reply(value)", "Respond to the caller immediately", COLOR_BLUE),
    ("andReply(persist(...), reply)", "Persist + reply atomically", COLOR_BLUE),
    ("andRun(persist(...), sideEffect)", "Persist, then run async side effect", COLOR_ORANGE),
    ("done()", "No-op — acknowledge without action", COLOR_LIGHT_GRAY),
]
for i, (code, desc, color) in enumerate(effects):
    y = Emu(2300000 + i * 700000)
    add_code_box(slide, code, M_LEFT, y, Emu(4800000), Emu(500000), font_size=15)
    add_text(slide, Emu(5900000), y + Emu(80000), Emu(5500000), Emu(400000),
             desc, font_size=17, color=color, bold=True)

add_notes(slide,
    "TIMING: 0:27–0:30\n"
    "Highlight persist + andReply — they'll use these in Exercise 1.\n"
    "Mention andRun — they'll see it in the LLM demo.\n"
    "\"decide returns a recipe, the runtime executes it.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11: The Effectful Event Loop
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_DARK)
add_text(slide, M_LEFT, Emu(200000), CONTENT_W, Emu(600000),
         "THE EFFECTFUL EVENT LOOP", font_size=30, bold=True, color=COLOR_WHITE)
add_accent_line(slide, Emu(800000), color=COLOR_BLUE)

# ── Layout constants ──
BOX_H = Emu(450000)
SMALL_H = Emu(380000)
ARROW_SZ = Emu(300000)

# ── Main loop (left column) ──────────────────────────────────────────────
# 5 steps down the left side, with side outputs

main_x = Emu(600000)
main_w = Emu(2600000)
side_x = Emu(3600000)    # where Journal/State' appear
side_w = Emu(1800000)

# Step 1: Command Queue
y1 = Emu(1100000)
add_box(slide, main_x, y1, main_w, BOX_H,
        COLOR_BLUE, "1  Dequeue command", font_size=14, text_color=COLOR_WHITE)
add_text(slide, main_x + main_w + Emu(100000), y1 + Emu(60000),
         Emu(2000000), SMALL_H,
         "from entity inbox", font_size=12, color=COLOR_LIGHT_GRAY)

# Arrow ↓
add_text(slide, main_x + Emu(1100000), y1 + BOX_H, Emu(400000), ARROW_SZ,
         "↓", font_size=22, bold=True, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Step 2: decide
y2 = Emu(1850000)
add_box(slide, main_x, y2, main_w, BOX_H,
        COLOR_DARK, "2  decide(state, cmd)", font_size=14, text_color=COLOR_GREEN)
add_text(slide, main_x + main_w + Emu(100000), y2 + Emu(60000),
         Emu(2200000), SMALL_H,
         "→ Effect<Event, Reply>", font_size=12, color=COLOR_GREEN, font_name="SF Mono")

# Arrow ↓
add_text(slide, main_x + Emu(1100000), y2 + BOX_H, Emu(400000), ARROW_SZ,
         "↓", font_size=22, bold=True, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Step 3: persist events
y3 = Emu(2600000)
add_box(slide, main_x, y3, main_w, BOX_H,
        RGBColor(0x1B, 0x5E, 0x20), "3  persist events", font_size=14, text_color=COLOR_WHITE)
# Side: → Journal
add_text(slide, main_x + main_w - Emu(100000), y3 + Emu(80000),
         Emu(500000), SMALL_H,
         "→", font_size=20, bold=True, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_box(slide, side_x, y3, side_w, BOX_H,
        RGBColor(0x1C, 0x19, 0x17), "Journal", font_size=13, text_color=COLOR_LIGHT_GRAY)

# Arrow ↓
add_text(slide, main_x + Emu(1100000), y3 + BOX_H, Emu(400000), ARROW_SZ,
         "↓", font_size=22, bold=True, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Step 4: apply events → State'
y4 = Emu(3350000)
add_box(slide, main_x, y4, main_w, BOX_H,
        RGBColor(0x1B, 0x5E, 0x20), "4  apply(state, event)", font_size=14, text_color=COLOR_WHITE)
# Side: → State'
add_text(slide, main_x + main_w - Emu(100000), y4 + Emu(80000),
         Emu(500000), SMALL_H,
         "→", font_size=20, bold=True, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_box(slide, side_x, y4, side_w, BOX_H,
        RGBColor(0x1C, 0x19, 0x17), "State'", font_size=13, text_color=COLOR_ORANGE)

# Arrow ↓
add_text(slide, main_x + Emu(1100000), y4 + BOX_H, Emu(400000), ARROW_SZ,
         "↓", font_size=22, bold=True, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Step 5: execute Effect
y5 = Emu(4100000)
add_box(slide, main_x, y5, main_w, BOX_H,
        COLOR_ORANGE, "5  execute Effect", font_size=14, text_color=COLOR_WHITE)

# ── Branch: reply (left-down) ────────────────────────────────────────────

reply_x = Emu(600000)
reply_y = Emu(4900000)
add_text(slide, reply_x + Emu(200000), y5 + BOX_H - Emu(30000),
         Emu(600000), ARROW_SZ,
         "↙", font_size=18, bold=True, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_box(slide, reply_x, reply_y, Emu(1600000), SMALL_H,
        COLOR_BLUE, "reply → Caller", font_size=12, text_color=COLOR_WHITE)

# ── Branch: andRun → ctx.sync (right side) ──────────────────────────────

sync_x = Emu(5900000)
sync_w = Emu(5200000)

# Arrow from step 5 to sync block
add_text(slide, main_x + main_w, y5 + Emu(80000),
         Emu(800000), SMALL_H,
         "andRun →", font_size=12, bold=True, color=COLOR_ORANGE, alignment=PP_ALIGN.CENTER)

# ctx.sync block (the main attraction)
sync_y = Emu(1100000)
sync_h = Emu(5200000)
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, sync_x, sync_y, sync_w, sync_h)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
shape.line.color.rgb = COLOR_ORANGE
shape.line.width = Pt(1.5)
shape.adjustments[0] = 0.02

add_text(slide, sync_x + Emu(100000), sync_y + Emu(80000),
         sync_w - Emu(200000), Emu(350000),
         "ctx.sync() — async effect execution", font_size=14, bold=True, color=COLOR_ORANGE)

# External service call
svc_y = Emu(1650000)
add_box(slide, sync_x + Emu(300000), svc_y, Emu(4600000), Emu(550000),
        RGBColor(0x2D, 0x2D, 0x2D), "", font_size=12)
add_text(slide, sync_x + Emu(400000), svc_y + Emu(50000),
         Emu(4400000), Emu(250000),
         "effect: () → call external service", font_size=13, color=COLOR_WHITE, font_name="SF Mono")
add_text(slide, sync_x + Emu(400000), svc_y + Emu(280000),
         Emu(4400000), Emu(250000),
         "LLM  •  payment gateway  •  shipping API  •  email", font_size=11, color=COLOR_LIGHT_GRAY)

# Three outcomes
outcome_y = Emu(2500000)
outcome_w = Emu(1400000)
outcome_gap = Emu(150000)
outcomes = [
    ("onSuccess", COLOR_GREEN, "result → Command"),
    ("onFailure", COLOR_RED, "error → Command"),
    ("onTimeout", COLOR_ORANGE, "→ Command"),
]
for i, (label, color, desc) in enumerate(outcomes):
    ox = sync_x + Emu(300000) + i * (outcome_w + outcome_gap)
    add_box(slide, ox, outcome_y, outcome_w, Emu(400000),
            color, label, font_size=12, text_color=COLOR_WHITE)
    add_text(slide, ox, outcome_y + Emu(420000), outcome_w, Emu(300000),
             desc, font_size=10, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow down from outcomes
add_text(slide, sync_x + Emu(2200000), Emu(3200000),
         Emu(800000), ARROW_SZ,
         "↓", font_size=22, bold=True, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# The command produced by sync
cmd_y = Emu(3550000)
add_box(slide, sync_x + Emu(800000), cmd_y, Emu(3600000), Emu(550000),
        COLOR_DARK, "", font_size=12)
add_text(slide, sync_x + Emu(900000), cmd_y + Emu(50000),
         Emu(3400000), Emu(250000),
         "→ new Command delivered to entity", font_size=13, bold=True, color=COLOR_WHITE)
add_text(slide, sync_x + Emu(900000), cmd_y + Emu(300000),
         Emu(3400000), Emu(250000),
         "e.g. { tag: \"SetEncouragement\", text: \"...\" }", font_size=11,
         color=COLOR_GREEN, font_name="SF Mono")

# ── Feedback arrow: Command → back to step 1 ────────────────────────────
# Visual: curved text showing the loop

loop_y = Emu(4400000)
add_text(slide, sync_x + Emu(800000), loop_y, Emu(3600000), Emu(600000),
         "↑  back to step 1 — the loop repeats\n"
         "   decide() handles the new command\n"
         "   → produces new events → new state",
         font_size=12, color=COLOR_ORANGE, bold=True)

# ── Vertical feedback arrow on the far right ─────────────────────────────
# A visible arrow from bottom to top along the right edge

arrow_x = sync_x + sync_w - Emu(300000)
add_text(slide, arrow_x, Emu(5100000), Emu(500000), Emu(500000),
         "⟲", font_size=36, bold=True, color=COLOR_ORANGE, alignment=PP_ALIGN.CENTER)

# ── Bottom label ─────────────────────────────────────────────────────────
add_text(slide, M_LEFT, Emu(6200000), CONTENT_W, Emu(400000),
         "One entity, one command at a time.  Effects are deterministic.  External calls get retry, replay, audit trail.",
         font_size=13, bold=True, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "TIMING: 0:30–0:33\n"
    "THIS IS THE KEY DIAGRAM. Walk through each step:\n"
    "1. \"Command arrives in the entity's inbox\"\n"
    "2. \"decide() looks at state + command, returns an Effect recipe\"\n"
    "3. \"Events get persisted to the journal — append-only, never mutate\"\n"
    "4. \"apply() folds events into new state — pure, deterministic\"\n"
    "5. \"The Effect executes: reply to caller, OR trigger an async side effect\"\n"
    "\"ctx.sync() is the magic — it calls any external service and routes the\n"
    "result back as a new command. Same loop, same guarantees.\"\n"
    "\"You'll see this live in the demo: Issue → mock service → SetEncouragement → new event.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12: DDD ↔ TEOB Mapping Table
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "THE DDD CORE MAPS 1:1")
add_accent_line(slide, Emu(1200000))

rows_data = [
    ("DDD", "TEOB", "Notes"),
    ("Aggregate", "Aggregate<S, C, E, R>", "The consistency boundary."),
    ("Aggregate ID", "EntityId", "Branded string, same concept."),
    ("Command", "Command (type param)", "Inbound intent."),
    ("Domain Event", "Event (type param)", "Fact that happened."),
    ("Decision function", "decide(state, cmd, ctx)", "\"Given state + command, what happens?\""),
    ("State fold / evolve", "apply(state, event)", "Pure left-fold."),
    ("Initial state", "initial(id)", "The \"zero\" of the fold."),
    ("Invariant", "invariants[]", "Executable, testable — faithful to DDD intent."),
]

col_widths = [Emu(2500000), Emu(3500000), Emu(5500000)]
table_shape = slide.shapes.add_table(
    len(rows_data), 3, M_LEFT, Emu(1500000),
    sum(col_widths), Emu(500000) * len(rows_data))
table = table_shape.table

for ci, w in enumerate(col_widths):
    table.columns[ci].width = w

for ri, row_data in enumerate(rows_data):
    for ci, cell_text in enumerate(row_data):
        cell = table.cell(ri, ci)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14 if ri > 0 else 15)
        p.font.bold = (ri == 0)
        if ri == 0:
            p.font.color.rgb = COLOR_WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BLUE
        elif ci == 1:
            p.font.name = "SF Mono"
            p.font.color.rgb = COLOR_BLUE
            p.font.size = Pt(13)
        else:
            p.font.color.rgb = COLOR_DARK if ci == 0 else COLOR_GRAY
        cell.margin_left = Pt(8)
        cell.margin_top = Pt(4)

add_notes(slide,
    "TIMING: 0:30–0:35\n"
    "THE AHA SLIDE. Go row by row.\n"
    "\"If you know DDD, you already know TEOB.\"\n"
    "Pause for questions before Exercise 1."
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12: Beyond the Blue Book
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "BEYOND THE BLUE BOOK")
add_accent_line(slide, Emu(1200000))

beyond = [
    ("Effect ADT", "Declarative Effect<Event, Reply> — a program describing what should happen."),
    ("Reply as first-class", "persist(...).andReply(...) atomically. CQRS-with-response."),
    ("EffectControl (ctx)", "Timers, cross-entity messaging, external effects (ctx.sync, ctx.tell)."),
    ("Invariants", "Executable and testable — faithful to why the boundary exists."),
    ("Projections", "Event → read model. Same pure fold pattern, optimized for queries."),
]
for i, (title, desc) in enumerate(beyond):
    y = Emu(1700000 + i * 900000)
    add_text(slide, M_LEFT, y, Emu(2800000), Emu(400000),
             title, font_size=18, bold=True, color=COLOR_BLUE)
    add_text(slide, Emu(3700000), y, Emu(8000000), Emu(700000),
             desc, font_size=15, color=COLOR_GRAY)

add_notes(slide,
    "TIMING: 0:35–0:38\n"
    "\"The core is DDD. The extras are what make it practical.\"\n"
    "Mention projections — they'll build one in Exercise 2."
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 13: Everything Is an Entity
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "EVERYTHING IS AN ENTITY")
add_accent_line(slide, Emu(1200000))

items = [
    "Complete entity behaviour — as pure functions",
    "Codecs — serialisable representation (Command, Reply, Event, State)",
    "Agnostic to execution and persistence",
    "Multiple runtimes: in-memory, SQLite, PostgreSQL (same aggregate code)",
]
add_bullet_block(slide, items, Emu(1700000), color=COLOR_DARK, font_size=19,
                 spacing=Emu(800000))

add_code_box(slide,
    "// Same aggregate code, different runtimes:\n"
    "createInMemoryRuntime([reg])     // tests & dev\n"
    "createSQLiteRuntime([reg])       // zero-config local\n"
    "createPostgresRuntime([reg])     // production",
    M_LEFT, Emu(5000000), Emu(7000000), Emu(1200000), font_size=15)

add_notes(slide,
    "TIMING: 0:38–0:40\n"
    "\"Your business logic doesn't know or care about the database.\"\n"
    "\"In exercises, we use InMemoryRuntime. In production, swap to Postgres.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 14: Event Storming — Discover the Domain
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "EVENT STORMING: DISCOVER THE DOMAIN", color=COLOR_ORANGE)
add_accent_line(slide, Emu(1200000), color=COLOR_ORANGE)

add_text(slide, M_LEFT, Emu(1500000), CONTENT_W, Emu(400000),
         "Before we code — what happened in the domain?",
         font_size=19, color=COLOR_GRAY)

# Sticky note vocabulary
COLOR_STICKY_ORANGE = RGBColor(0xFF, 0x9E, 0x22)
COLOR_STICKY_BLUE = RGBColor(0x42, 0xA5, 0xF5)
COLOR_STICKY_YELLOW = RGBColor(0xFF, 0xEE, 0x58)
COLOR_STICKY_PINK = RGBColor(0xEF, 0x53, 0x50)
COLOR_STICKY_GREEN = RGBColor(0x66, 0xBB, 0x6A)

stickies = [
    (COLOR_STICKY_ORANGE, "Domain Event", "Something that happened\n(past tense)", COLOR_DARK),
    (COLOR_STICKY_BLUE, "Command", "User intent\n(imperative)", COLOR_WHITE),
    (COLOR_STICKY_YELLOW, "Aggregate", "Consistency boundary\n(the decision maker)", COLOR_DARK),
    (COLOR_STICKY_PINK, "Business Rule", "Invariant / policy\n(must always hold)", COLOR_WHITE),
    (COLOR_STICKY_GREEN, "Read Model", "Query-optimized view\n(derived from events)", COLOR_WHITE),
]
sticky_w = Emu(2000000)
sticky_h = Emu(1600000)
sticky_gap = Emu(150000)
total_sticky_w = len(stickies) * sticky_w + (len(stickies) - 1) * sticky_gap
sticky_start_x = Emu((SLIDE_W - total_sticky_w) // 2)

for i, (color, label, desc, text_color) in enumerate(stickies):
    x = sticky_start_x + i * (sticky_w + sticky_gap)
    y = Emu(2200000)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, sticky_w, sticky_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = -2 + (i % 3) * 2  # slight tilt for sticky note feel
    # Label
    add_text(slide, x + Emu(80000), y + Emu(80000), sticky_w - Emu(160000), Emu(400000),
             label, font_size=14, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
    # Description
    add_text(slide, x + Emu(80000), y + Emu(550000), sticky_w - Emu(160000), Emu(900000),
             desc, font_size=11, color=text_color, alignment=PP_ALIGN.CENTER)

add_text(slide, M_LEFT, Emu(4300000), CONTENT_W, Emu(600000),
         "Alberto Brandolini's collaborative modelling technique — sticky notes on a wall",
         font_size=15, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text(slide, M_LEFT, Emu(5000000), CONTENT_W, Emu(600000),
         "\"Let's storm the gift card domain together...\"",
         font_size=22, bold=True, color=COLOR_ORANGE, alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "TIMING: 0:40–0:45\n"
    "\"Before we write code, DDD says we discover the domain first.\"\n"
    "\"Event Storming is how — sticky notes, a wall, and a conversation.\"\n"
    "Point to each colour. Ask the room: \"What events happen to a gift card?\"\n"
    "Let them call out events before showing the next slide."
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 15: Gift Card Event Storming Result
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "GIFT CARD: FROM STICKY NOTES TO CODE", color=COLOR_ORANGE)
add_accent_line(slide, Emu(1200000), color=COLOR_ORANGE)

# Left side: Event Storming sticky notes layout
es_left = M_LEFT

# Commands (blue) column
add_text(slide, es_left, Emu(1500000), Emu(2000000), Emu(300000),
         "COMMANDS", font_size=12, bold=True, color=COLOR_LIGHT_GRAY)
cmds = ["Issue", "Redeem", "Cancel", "GetBalance"]
for i, cmd in enumerate(cmds):
    y = Emu(1850000 + i * 520000)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, es_left, y, Emu(1800000), Emu(420000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_STICKY_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = cmd
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

# Arrow
add_text(slide, Emu(2200000), Emu(2700000), Emu(800000), Emu(400000),
         "→", font_size=36, bold=True, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

# Aggregate (yellow) — center
agg_x = Emu(3100000)
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, agg_x, Emu(2000000), Emu(2200000), Emu(2000000))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_STICKY_YELLOW
shape.line.fill.background()
add_text(slide, agg_x + Emu(100000), Emu(2100000), Emu(2000000), Emu(400000),
         "GiftCard", font_size=22, bold=True, color=COLOR_DARK, alignment=PP_ALIGN.CENTER)
# Invariants inside aggregate
inv_items = ["balance ≥ 0", "cancelled → bal = 0"]
for i, inv in enumerate(inv_items):
    add_text(slide, agg_x + Emu(100000), Emu(2700000 + i * 400000),
             Emu(2000000), Emu(350000),
             f"⚠ {inv}", font_size=12, color=COLOR_STICKY_PINK, alignment=PP_ALIGN.CENTER)

# Arrow
add_text(slide, Emu(5400000), Emu(2700000), Emu(800000), Emu(400000),
         "→", font_size=36, bold=True, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

# Events (orange) column
evt_x = Emu(6300000)
add_text(slide, evt_x, Emu(1500000), Emu(2000000), Emu(300000),
         "DOMAIN EVENTS", font_size=12, bold=True, color=COLOR_LIGHT_GRAY)
evts = ["Issued", "Redeemed", "Cancelled"]
for i, evt in enumerate(evts):
    y = Emu(1850000 + i * 520000)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, evt_x, y, Emu(1800000), Emu(420000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_STICKY_ORANGE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = evt
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK
    p.alignment = PP_ALIGN.CENTER

# Read Model (green) — far right
rm_x = Emu(8700000)
add_text(slide, rm_x, Emu(1500000), Emu(2800000), Emu(300000),
         "READ MODEL", font_size=12, bold=True, color=COLOR_LIGHT_GRAY)
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, rm_x, Emu(1850000), Emu(2500000), Emu(2000000))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_STICKY_GREEN
shape.line.fill.background()
rm_fields = ["balance", "status", "recipientName", "transactionCount"]
for i, field in enumerate(rm_fields):
    add_text(slide, rm_x + Emu(100000), Emu(2000000 + i * 380000),
             Emu(2300000), Emu(300000),
             field, font_size=13, color=COLOR_WHITE, font_name="SF Mono",
             alignment=PP_ALIGN.CENTER)

# Bottom: code mapping
add_code_box(slide,
    "// Sticky notes → TypeScript types:\n"
    "Command:   Issue | Redeem | Cancel | GetBalance\n"
    "Event:     Issued | Redeemed | Cancelled\n"
    "State:     { balance, status, recipientName }\n"
    "View:      { balance, status, recipientName, transactionCount }",
    M_LEFT, Emu(4500000), Emu(8500000), Emu(1500000), font_size=14)

add_text(slide, Emu(9200000), Emu(4700000), Emu(2500000), Emu(400000),
         "→ This is what\n   you'll implement!",
         font_size=18, bold=True, color=COLOR_GREEN)

add_notes(slide,
    "TIMING: 0:45–0:50\n"
    "Walk through the board left to right:\n"
    "1. Commands = user intent (blue)\n"
    "2. Aggregate = decision maker + rules (yellow)\n"
    "3. Events = facts that happened (orange)\n"
    "4. Read Model = query view (green)\n"
    "\"Every sticky note maps directly to a type in your code.\"\n"
    "\"Now let's implement it.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 16: Exercise 1 Brief — Gift Card Aggregate
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "EXERCISE 1: GIFT CARD AGGREGATE", color=COLOR_GREEN, font_size=34)
add_accent_line(slide, Emu(1200000), color=COLOR_GREEN)

add_code_box(slide,
    "Commands:  Issue | Redeem | Cancel | GetBalance\n"
    "Events:    Issued | Redeemed | Cancelled\n"
    "State:     { balance, status, recipientName }",
    M_LEFT, Emu(1600000), Emu(6500000), Emu(1000000), font_size=17)

add_text(slide, M_LEFT, Emu(2900000), Emu(6000000), Emu(400000),
         "Business rules to implement:", font_size=18, bold=True, color=COLOR_DARK)
rules = [
    "Issue only if card is Empty",
    "Redeem only if Active, amount ≤ balance",
    "Cancel only if Active",
    "Invariants: balance ≥ 0, cancelled → balance = 0",
]
add_bullet_block(slide, rules, Emu(3400000), color=COLOR_DARK, font_size=17,
                 spacing=Emu(420000))

add_code_box(slide,
    "npm run test:aggregate",
    M_LEFT, Emu(5400000), Emu(5500000), Emu(500000), font_size=16)

add_text(slide, Emu(6500000), Emu(5450000), Emu(5000000), Emu(500000),
         "16 tests. Make them green.", font_size=22, bold=True, color=COLOR_GREEN)

add_notes(slide,
    "TIMING: 0:50–0:52\n"
    "LIVE DEMO: Run tests. Show 14 failing, 2 passing.\n"
    "\"Start with Issue — it's the simplest. The hints are in HINTS-aggregate.md.\"\n"
    "\"Solution is in aggregate.solution.ts — but try first!\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 15: Exercise 1 — Work Time
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_DARK)
add_text(slide, M_LEFT, Emu(2200000), CONTENT_W, Emu(1000000),
         "EXERCISE 1: AGGREGATE", font_size=48, bold=True,
         color=COLOR_GREEN, alignment=PP_ALIGN.CENTER)
add_text(slide, M_LEFT, Emu(3400000), CONTENT_W, Emu(600000),
         "30 minutes  •  16 tests  •  implement decide() + invariants",
         font_size=22, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, M_LEFT, Emu(4600000), CONTENT_W, Emu(400000),
         "npm run test:aggregate",
         font_size=18, color=COLOR_GREEN, alignment=PP_ALIGN.CENTER,
         font_name="SF Mono")

add_notes(slide,
    "TIMING: 0:52–1:22 (30 min)\n"
    "LEAVE THIS SLIDE UP.\n"
    "At 15 min: \"How many tests green? More than 5?\"\n"
    "At 25 min: Live-code Issue case if people are stuck.\n"
    "At 28 min: \"2 minutes left.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 16: Exercise 1 Debrief
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "EXERCISE 1: DEBRIEF")
add_accent_line(slide, Emu(1200000))

takeaways = [
    "decide() — all business logic in one place, pattern-matched on command tag",
    "apply() — pure state fold, no logic, just data transformation",
    "invariants[] — executable contracts that guard consistency",
    "Effect ADT — persist(), reply(), andReply() — declarative, composable",
    "AggregateTestKit — test the domain directly, no infrastructure",
]
add_bullet_block(slide, takeaways, Emu(1700000), color=COLOR_DARK, font_size=17,
                 spacing=Emu(550000))

add_code_box(slide,
    "// The core pattern:\n"
    "if (invalid) return reply({ tag: 'Rejected', reason: '...' })\n"
    "return andReply(persist(event), { tag: 'Ok' })",
    M_LEFT, Emu(5200000), Emu(8500000), Emu(1000000), font_size=15)

add_notes(slide,
    "TIMING: 1:22–1:27\n"
    "Live-code the solution if needed.\n"
    "\"No database, no HTTP, no framework setup. Just domain logic.\"\n"
    "Ask: \"What surprised you?\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 17: Demo — HTTP Service
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "DEMO: YOUR AGGREGATE AS AN API", color=COLOR_BLUE)
add_accent_line(slide, Emu(1200000), color=COLOR_BLUE)

add_text(slide, M_LEFT, Emu(1600000), CONTENT_W, Emu(400000),
         "Your aggregate is now a REST API — defined by openapi.yaml:",
         font_size=19, color=COLOR_GRAY)

add_code_box(slide,
    "# Issue a gift card\n"
    "curl -X POST http://localhost:3000/api/gift-cards \\\n"
    "  -H 'Content-Type: application/json' \\\n"
    "  -d '{\"id\":\"card-1\",\"amount\":100,\"recipientName\":\"Alice\"}'\n"
    "→ 201 { balance: 100, status: \"Active\", recipientName: \"Alice\",\n"
    "         encouragement: \"Alice, you deserve something special!\" }\n\n"
    "# Redeem\n"
    "curl -X POST .../api/gift-cards/card-1/redeem -d '{\"amount\":30}'\n"
    "→ 200 { balance: 70, transactionCount: 1, encouragement: \"...\" }\n\n"
    "# Reject (too much)\n"
    "curl -X POST .../api/gift-cards/card-1/redeem -d '{\"amount\":999}'\n"
    "→ 400 { error: \"Insufficient balance\" }",
    M_LEFT, Emu(2200000), Emu(8500000), Emu(3400000), font_size=14)

add_text(slide, M_LEFT, Emu(5700000), CONTENT_W, Emu(400000),
         "REST endpoints return views, not raw replies.  Commands stay internal.  OpenAPI spec = contract.",
         font_size=15, bold=True, color=COLOR_BLUE)

add_notes(slide,
    "TIMING: 1:28–1:31\n"
    "LIVE DEMO: Run `npm start`, then curl commands.\n"
    "Show: Issue returns view WITH encouragement.\n"
    "Show: Redeem updates balance and transactionCount.\n"
    "Show: rejection → 400 with error message.\n"
    "\"Wait — where did the encouragement come from? Let me show you...\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 18: What Just Happened — Issue Flow with ctx.sync
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_DARK)
add_text(slide, M_LEFT, Emu(200000), CONTENT_W, Emu(600000),
         "WHAT JUST HAPPENED?", font_size=30, bold=True, color=COLOR_WHITE)
add_accent_line(slide, Emu(800000), color=COLOR_ORANGE)

# ── Sequence diagram ─────────────────────────────────────────────────────
# Participants across the top
parts = [
    ("Client", Emu(300000), COLOR_BLUE),
    ("HTTP", Emu(2200000), COLOR_BLUE),
    ("decide()", Emu(4100000), COLOR_GREEN),
    ("Journal", Emu(6000000), RGBColor(0x1C, 0x19, 0x17)),
    ("Mock Service", Emu(7900000), COLOR_ORANGE),
    ("Projection", Emu(9800000), RGBColor(0x17, 0x25, 0x54)),
]
part_w = Emu(1600000)
part_h = Emu(380000)
part_y = Emu(1050000)

for label, px, color in parts:
    add_box(slide, px, part_y, part_w, part_h, color,
            text=label, font_size=11, text_color=COLOR_WHITE)

# ── Lifelines (vertical dashed look via thin boxes) ──────────────────────
line_top = part_y + part_h
line_h = Emu(4700000)
for _, px, color in parts:
    cx = px + part_w // 2 - Emu(10000)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, cx, line_top, Emu(20000), line_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x44, 0x44, 0x44)
    shape.line.fill.background()

# ── Messages (horizontal arrows as text) ─────────────────────────────────
# Each message: (from_idx, to_idx, y_offset, label, color)
msg_font = 10
row_h = Emu(420000)

def msg_arrow(from_idx, to_idx, row, label, color=COLOR_WHITE, style="→"):
    from_x = parts[from_idx][1] + part_w // 2
    to_x = parts[to_idx][1] + part_w // 2
    y = line_top + Emu(80000) + row * row_h
    left = min(from_x, to_x) + Emu(30000)
    width = abs(to_x - from_x) - Emu(60000)
    # Arrow line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, y + Emu(180000), width, Emu(15000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Arrowhead (triangle)
    if to_idx > from_idx:
        ah_x = left + width - Emu(10000)
    else:
        ah_x = left - Emu(100000)
    arrow_dir = "→" if to_idx > from_idx else "←"
    add_text(slide, ah_x, y + Emu(100000), Emu(200000), Emu(200000),
             arrow_dir, font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # Label above
    add_text(slide, left, y, width, Emu(180000),
             label, font_size=msg_font, color=color, alignment=PP_ALIGN.CENTER)

# Row 0: Client → HTTP: POST /api/gift-cards
msg_arrow(0, 1, 0, "POST /api/gift-cards", COLOR_BLUE)

# Row 1: HTTP → decide: Issue(100, "Alice")
msg_arrow(1, 2, 1, "Issue(100, \"Alice\")", COLOR_GREEN)

# Row 2: decide → Journal: persist(Issued)
msg_arrow(2, 3, 2, "persist(Issued)", COLOR_GREEN)

# Row 3: decide → Mock Service: andRun → ctx.sync()
msg_arrow(2, 4, 3, "andRun → ctx.sync()", COLOR_ORANGE)

# Row 4: Mock Service → decide: SetEncouragement("...")
msg_arrow(4, 2, 4, "SetEncouragement(text)", COLOR_ORANGE)

# Row 5: decide → Journal: persist(EncouragementSet)
msg_arrow(2, 3, 5, "persist(EncouragementSet)", COLOR_ORANGE)

# Row 6: HTTP → Projection: refresh
msg_arrow(1, 5, 6, "refreshProjection()", RGBColor(0x88, 0xBB, 0xFF))

# Row 7: Projection → HTTP → Client: 201 { ..., encouragement }
msg_arrow(5, 1, 7, "GiftCardView", RGBColor(0x88, 0xBB, 0xFF))

# Row 8: HTTP → Client: response
msg_arrow(1, 0, 8.5, "201 { balance, encouragement, ... }", COLOR_BLUE)

# Row 9 (annotation): highlight the loop
ann_x = Emu(6200000)
ann_y = line_top + Emu(80000) + int(4.5 * row_h)
add_text(slide, ann_x + Emu(2200000), ann_y - Emu(50000), Emu(3200000), Emu(500000),
         "⟲ back through decide()\n   same loop, new event",
         font_size=9, bold=True, color=COLOR_ORANGE)

# ── Bottom annotation ────────────────────────────────────────────────────
add_text(slide, M_LEFT, Emu(6200000), CONTENT_W, Emu(400000),
         "ctx.sync() result → Command → decide() → Event.  The loop from slide 11. Same pattern for LLM, payment, anything.",
         font_size=12, bold=True, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "TIMING: 1:31–1:35\n"
    "\"Remember the event loop from earlier? This is what just happened.\"\n"
    "Walk through each arrow:\n"
    "1. Client POST → HTTP layer maps to Issue command\n"
    "2. decide() validates, returns persist(Issued) + andRun\n"
    "3. Issued event saved to journal\n"
    "4. andRun fires ctx.sync → mock service generates encouragement\n"
    "5. Result comes back as SetEncouragement COMMAND — back through decide()!\n"
    "6. decide() handles SetEncouragement → persist(EncouragementSet)\n"
    "7. Projection rebuilds view from both events\n"
    "8. HTTP returns the complete view with encouragement\n"
    "\"Same pattern works for real LLM, payment gateway, shipping API.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 19: Break
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_DARK)
add_text(slide, M_LEFT, Emu(2600000), CONTENT_W, Emu(800000),
         "BREAK", font_size=64, bold=True,
         color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
add_text(slide, M_LEFT, Emu(3600000), CONTENT_W, Emu(600000),
         "5 minutes  •  stretch  •  refill coffee",
         font_size=22, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, M_LEFT, Emu(4800000), CONTENT_W, Emu(400000),
         "Next up: the read side — projections",
         font_size=18, color=COLOR_GREEN, alignment=PP_ALIGN.CENTER)

add_notes(slide, "TIMING: 1:32–1:37")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 19: Theory — Write Side vs Read Side
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "WRITE SIDE VS READ SIDE")
add_accent_line(slide, Emu(1200000))

# Left: write side
add_box(slide, M_LEFT, Emu(1800000), Emu(4800000), Emu(2500000),
        COLOR_BLUE, "", font_size=14)
add_text(slide, Emu(900000), Emu(1900000), Emu(4400000), Emu(400000),
         "WRITE SIDE (aggregate)", font_size=18, bold=True, color=COLOR_WHITE)
write_items = [
    "decide() — validates, produces events",
    "apply() — state for decisions",
    "Optimized for consistency",
    "One entity at a time",
]
for i, item in enumerate(write_items):
    add_text(slide, Emu(900000), Emu(2400000 + i * 400000), Emu(4400000), Emu(350000),
             f"• {item}", font_size=14, color=COLOR_WHITE)

# Right: read side
add_box(slide, Emu(6200000), Emu(1800000), Emu(4800000), Emu(2500000),
        COLOR_GREEN, "", font_size=14)
add_text(slide, Emu(6400000), Emu(1900000), Emu(4400000), Emu(400000),
         "READ SIDE (projection)", font_size=18, bold=True, color=COLOR_WHITE)
read_items = [
    "evolve() — builds query-optimized views",
    "Can include derived data (counts, aggregations)",
    "Optimized for queries",
    "Can span multiple entities",
]
for i, item in enumerate(read_items):
    add_text(slide, Emu(6400000), Emu(2400000 + i * 400000), Emu(4400000), Emu(350000),
             f"• {item}", font_size=14, color=COLOR_WHITE)

# Arrow
add_text(slide, Emu(5100000), Emu(2700000), Emu(1000000), Emu(400000),
         "→", font_size=36, bold=True, color=COLOR_ORANGE, alignment=PP_ALIGN.CENTER)

add_text(slide, M_LEFT, Emu(4800000), CONTENT_W, Emu(400000),
         "Same pattern: pure function, event in → new state out",
         font_size=20, bold=True, color=COLOR_DARK, alignment=PP_ALIGN.CENTER)

add_code_box(slide,
    "// Write side                          // Read side\n"
    "apply(state, event) → State           evolve(view, event) → View",
    M_LEFT, Emu(5400000), CONTENT_W, Emu(600000), font_size=16)

add_notes(slide,
    "TIMING: 1:37–1:42\n"
    "\"You already know the pattern — apply and evolve are the same shape.\"\n"
    "\"The difference: evolve builds views optimized for queries, not decisions.\"\n"
    "\"A view can include transactionCount — something the aggregate doesn't track.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 20: Exercise 2 Brief — Projection
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "EXERCISE 2: GIFT CARD PROJECTION", color=COLOR_GREEN, font_size=34)
add_accent_line(slide, Emu(1200000), color=COLOR_GREEN)

add_text(slide, M_LEFT, Emu(1600000), CONTENT_W, Emu(400000),
         "Build a read model from gift card events:", font_size=19, color=COLOR_GRAY)

add_code_box(slide,
    "GiftCardView {\n"
    "  balance, status, recipientName,\n"
    "  encouragement,        // from LLM (demo)\n"
    "  transactionCount      // derived — aggregate doesn't track this!\n"
    "}",
    M_LEFT, Emu(2200000), Emu(6500000), Emu(1400000), font_size=16)

add_text(slide, M_LEFT, Emu(3900000), Emu(6000000), Emu(400000),
         "Implement evolve() — same pattern as apply():", font_size=18, bold=True, color=COLOR_DARK)
cases = [
    "Issued → set balance, status, recipientName",
    "Redeemed → subtract amount, increment transactionCount",
    "Cancelled → zero balance, set status",
    "EncouragementSet → set encouragement text",
]
add_bullet_block(slide, cases, Emu(4400000), color=COLOR_DARK, font_size=16,
                 spacing=Emu(380000))

add_code_box(slide,
    "npm run test:projection",
    M_LEFT, Emu(5900000), Emu(5500000), Emu(500000), font_size=16)

add_text(slide, Emu(6500000), Emu(5950000), Emu(5000000), Emu(500000),
         "8 tests. Make them green.", font_size=22, bold=True, color=COLOR_GREEN)

add_notes(slide,
    "TIMING: 1:42–1:45\n"
    "\"You already know the pattern from apply(). This is the same thing, for queries.\"\n"
    "\"The key difference: transactionCount. The aggregate doesn't track this — the projection does.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 21: Exercise 2 — Work Time
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_DARK)
add_text(slide, M_LEFT, Emu(2200000), CONTENT_W, Emu(1000000),
         "EXERCISE 2: PROJECTION", font_size=48, bold=True,
         color=COLOR_GREEN, alignment=PP_ALIGN.CENTER)
add_text(slide, M_LEFT, Emu(3400000), CONTENT_W, Emu(600000),
         "20 minutes  •  8 tests  •  implement evolve()",
         font_size=22, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, M_LEFT, Emu(4600000), CONTENT_W, Emu(400000),
         "npm run test:projection",
         font_size=18, color=COLOR_GREEN, alignment=PP_ALIGN.CENTER,
         font_name="SF Mono")
add_text(slide, M_LEFT, Emu(5400000), CONTENT_W, Emu(400000),
         "Same pattern as apply()  •  HINTS-projection.md if stuck",
         font_size=16, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "TIMING: 1:45–2:05 (20 min)\n"
    "This should go faster — they already know the pattern.\n"
    "At 10 min: Most should have Issued done.\n"
    "At 15 min: Show Hint 3 (Redeemed with transactionCount).\n"
    "At 18 min: \"2 minutes left.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 22: Exercise 2 Debrief
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "EXERCISE 2: DEBRIEF")
add_accent_line(slide, Emu(1200000))

add_text(slide, M_LEFT, Emu(1600000), CONTENT_W, Emu(400000),
         "Three pure functions — that's the whole model:", font_size=20, bold=True, color=COLOR_DARK)

add_code_box(slide,
    "decide(state, command)  → Effect     // business logic\n"
    "apply(state, event)     → State      // write-side fold\n"
    "evolve(view, event)     → View       // read-side fold",
    M_LEFT, Emu(2200000), Emu(8000000), Emu(1000000), font_size=18)

add_text(slide, M_LEFT, Emu(3600000), CONTENT_W, Emu(400000),
         "What the projection gave you that the aggregate doesn't:", font_size=18, bold=True, color=COLOR_DARK)

diffs = [
    "transactionCount — derived data, computed from event stream",
    "encouragement — data from a different event (LLM callback)",
    "Optimized for queries — flat shape, no business rules",
    "Can be rebuilt anytime from the event journal",
]
add_bullet_block(slide, diffs, Emu(4100000), color=COLOR_BLUE, font_size=17,
                 spacing=Emu(500000))

add_notes(slide,
    "TIMING: 2:05–2:10\n"
    "\"Three pure functions. That's the entire architecture.\"\n"
    "\"decide for commands, apply for state, evolve for queries.\"\n"
    "\"Now something fun — let's get a frontend in 10 minutes.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 25: Exercise 3 Brief — AI-Generate a Frontend
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "EXERCISE 3: AI-GENERATE A FRONTEND", color=COLOR_GREEN, font_size=32)
add_accent_line(slide, Emu(1200000), color=COLOR_GREEN)

add_text(slide, M_LEFT, Emu(1600000), CONTENT_W, Emu(400000),
         "Your service has both write and read APIs. One prompt → working UI.",
         font_size=19, color=COLOR_GRAY)

add_text(slide, M_LEFT, Emu(2200000), CONTENT_W, Emu(400000),
         "The setup:", font_size=18, bold=True, color=COLOR_DARK)

steps = [
    "Open PROMPT-frontend.md — it references openapi.yaml",
    "Copy the prompt into your AI coding agent (Claude Code, Cursor, etc.)",
    "It generates public/index.html — vanilla HTML/CSS/JS, no build tools",
    "npm start → open http://localhost:3000",
]
add_bullet_block(slide, steps, Emu(2700000), color=COLOR_DARK, font_size=17,
                 spacing=Emu(500000))

add_text(slide, M_LEFT, Emu(5000000), CONTENT_W, Emu(500000),
         "The AI reads the OpenAPI spec and writes the frontend for you.\n"
         "Reference solution: public-solution/index.html",
         font_size=16, color=COLOR_LIGHT_GRAY)

add_notes(slide,
    "TIMING: 2:10–2:12\n"
    "\"You've built the backend — aggregate + projection. Now let's get a UI.\"\n"
    "\"The prompt references openapi.yaml — the AI gets type safety from the spec.\"\n"
    "Show PROMPT-frontend.md briefly."
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 26: Exercise 3 — Work Time
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_DARK)
add_text(slide, M_LEFT, Emu(2200000), CONTENT_W, Emu(1000000),
         "EXERCISE 3: AI FRONTEND", font_size=48, bold=True,
         color=COLOR_GREEN, alignment=PP_ALIGN.CENTER)
add_text(slide, M_LEFT, Emu(3400000), CONTENT_W, Emu(600000),
         "10 minutes  •  one prompt  •  full working UI",
         font_size=22, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, M_LEFT, Emu(4400000), CONTENT_W, Emu(400000),
         "PROMPT-frontend.md → your AI agent → public/index.html",
         font_size=18, color=COLOR_GREEN, alignment=PP_ALIGN.CENTER,
         font_name="SF Mono")
add_text(slide, M_LEFT, Emu(5400000), CONTENT_W, Emu(400000),
         "npm start  •  http://localhost:3000",
         font_size=16, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "TIMING: 2:12–2:22 (10 min)\n"
    "LEAVE THIS SLIDE UP.\n"
    "Walk around — help anyone whose agent produces broken code.\n"
    "At 5 min: Show the reference solution if people want to compare.\n"
    "At 8 min: \"2 minutes left — make sure npm start works.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 27: Demo — LLM Integration
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "DEMO: LLM INTEGRATION", color=COLOR_PURPLE)
add_accent_line(slide, Emu(1200000), color=COLOR_PURPLE)

add_text(slide, M_LEFT, Emu(1600000), CONTENT_W, Emu(400000),
         "When a gift card is issued, ask an LLM to write a personal encouragement:",
         font_size=18, color=COLOR_GRAY)

add_code_box(slide,
    "case \"Issue\": {\n"
    "  return andReply(\n"
    "    andRun(\n"
    "      persist({ tag: \"Issued\", ... }),\n"
    "      async () => {\n"
    "        await ctx.sync({                        // ← same pattern!\n"
    "          effect: () => generateEncouragement(name, amount),\n"
    "          onSuccess: (r) => ({ tag: \"SetEncouragement\", text: r.text }),\n"
    "          onFailure: (_) => ({ tag: \"SetEncouragement\", text: fallback }),\n"
    "        });\n"
    "      },\n"
    "    ),\n"
    "    { tag: \"Ok\" },\n"
    "  );\n"
    "}",
    M_LEFT, Emu(2200000), Emu(9000000), Emu(3200000), font_size=14)

add_text(slide, M_LEFT, Emu(5700000), CONTENT_W, Emu(400000),
         "ctx.sync() — same pattern for LLM, payment gateway, shipping API. Event sourcing gives you retry, replay, audit — for free.",
         font_size=16, bold=True, color=COLOR_PURPLE)

add_notes(slide,
    "TIMING: 2:22–2:32\n"
    "LIVE DEMO: Run `OPENROUTER_API_KEY=sk-... npm run demo:llm`\n"
    "Issue a card with a name → show the encouragement appearing.\n"
    "\"ctx.sync() is the same pattern you'd use for ANY external call.\"\n"
    "\"The LLM result comes back as a command, gets persisted as an event.\"\n"
    "\"Crash? Replay from journal. Debug? Replay specific steps.\"\n"
    "THIS IS THE WOW MOMENT."
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 28: The Full Picture — Architecture
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "THE FULL PICTURE")
add_accent_line(slide, Emu(1200000))

# ── Browser ──
add_box(slide, Emu(4500000), Emu(1500000), Emu(3000000), Emu(500000),
        COLOR_LIGHT_GRAY, "Browser — localhost:3000", font_size=13)

# Arrow down
add_text(slide, Emu(5600000), Emu(2050000), Emu(800000), Emu(400000),
         "↓", font_size=24, bold=True, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

# ── HTTP Layer ──
add_box(slide, Emu(1500000), Emu(2400000), Emu(9000000), Emu(600000),
        COLOR_BLUE, "", font_size=12)
add_text(slide, Emu(1700000), Emu(2450000), Emu(8600000), Emu(250000),
         "HTTP Layer — Hono (REST → TEOB commands)", font_size=13, bold=True,
         color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
endpoints = "POST /gift-cards   GET /gift-cards   GET /:id   POST /:id/redeem   POST /:id/cancel"
add_text(slide, Emu(1700000), Emu(2700000), Emu(8600000), Emu(250000),
         endpoints, font_size=10, color=RGBColor(0xBB, 0xDD, 0xFF),
         font_name="SF Mono", alignment=PP_ALIGN.CENTER)

# Arrow down to write and read sides
add_text(slide, Emu(3500000), Emu(3050000), Emu(800000), Emu(400000),
         "↓", font_size=24, bold=True, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, Emu(7700000), Emu(3050000), Emu(800000), Emu(400000),
         "↓", font_size=24, bold=True, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

# ── Write Side ──
ws_x = Emu(1500000)
ws_y = Emu(3400000)
ws_w = Emu(4200000)
ws_h = Emu(2400000)
add_box(slide, ws_x, ws_y, ws_w, ws_h, RGBColor(0x05, 0x2E, 0x16), "", font_size=12)
add_text(slide, ws_x + Emu(100000), ws_y + Emu(80000), ws_w - Emu(200000), Emu(300000),
         "WRITE SIDE — Aggregate", font_size=14, bold=True, color=COLOR_GREEN)

add_code_box(slide,
    "decide(state, cmd) → Effect\n"
    "apply(state, event) → State",
    ws_x + Emu(150000), ws_y + Emu(500000), ws_w - Emu(300000), Emu(650000), font_size=12)

add_text(slide, ws_x + Emu(150000), ws_y + Emu(1300000), ws_w - Emu(300000), Emu(800000),
         "Invariants:\n• balance ≥ 0\n• cancelled → balance = 0",
         font_size=11, color=COLOR_GREEN)

# Journal
j_x = ws_x + Emu(150000)
j_y = ws_y + ws_h + Emu(200000)
add_box(slide, j_x, j_y, ws_w - Emu(300000), Emu(400000),
        RGBColor(0x1C, 0x19, 0x17), "Journal (events per entity)", font_size=11,
        text_color=COLOR_LIGHT_GRAY)

# ── Read Side ──
rs_x = Emu(6300000)
rs_y = Emu(3400000)
rs_w = Emu(4200000)
rs_h = Emu(2400000)
add_box(slide, rs_x, rs_y, rs_w, rs_h, RGBColor(0x17, 0x25, 0x54), "", font_size=12)
add_text(slide, rs_x + Emu(100000), rs_y + Emu(80000), rs_w - Emu(200000), Emu(300000),
         "READ SIDE — Projection", font_size=14, bold=True, color=COLOR_WHITE)

add_code_box(slide,
    "evolve(view, event) → View",
    rs_x + Emu(150000), rs_y + Emu(500000), rs_w - Emu(300000), Emu(450000), font_size=12)

add_text(slide, rs_x + Emu(150000), rs_y + Emu(1100000), rs_w - Emu(300000), Emu(1100000),
         "GiftCardView:\n  balance, status,\n  recipientName,\n  encouragement,\n  transactionCount",
         font_size=11, color=RGBColor(0x88, 0xBB, 0xFF), font_name="SF Mono")

# Projection Store
ps_x = rs_x + Emu(150000)
ps_y = rs_y + rs_h + Emu(200000)
add_box(slide, ps_x, ps_y, rs_w - Emu(300000), Emu(400000),
        RGBColor(0x1C, 0x19, 0x17), "Projection Store (views per entity)", font_size=11,
        text_color=COLOR_LIGHT_GRAY)

# Arrow: Journal → Read Side (events flow)
add_text(slide, Emu(5000000), Emu(4500000), Emu(1300000), Emu(400000),
         "events →", font_size=12, bold=True, color=COLOR_ORANGE, alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "TIMING: 2:32–2:35\n"
    "\"You built the green layer (aggregate) and the blue layer (projection).\"\n"
    "\"The HTTP layer is a thin REST mapping — commands stay internal.\"\n"
    "\"To go to production: swap InMemoryRuntime → PostgresRuntime.\""
)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 29: Battle Tested
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_LIGHT)
add_title(slide, "BATTLE TESTED")
add_accent_line(slide, Emu(1200000))

points = [
    "Natively CQS — for (almost) infinite scaling",
    "FP and DDD — to express intrinsically complex data and rules",
    "Everything is executed deterministically",
    "10+ years of real production (fintech) experience",
]
add_bullet_block(slide, points, Emu(1800000), color=COLOR_DARK, font_size=21,
                 spacing=Emu(900000), bold=True)

add_notes(slide, "TIMING: 2:35–2:37\nCredibility slide. Keep brief.")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 30: Closing CTA
# ══════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(slide, BG_DARK)
add_text(slide, M_LEFT, Emu(1000000), CONTENT_W, Emu(800000),
         "WHAT'S NEXT?", font_size=44, bold=True,
         color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Emu(1800000), color=COLOR_BLUE)

cta_items = [
    ("Clone the repo", "github.com/lambda-house/teob-ts-workshop-gift-card", COLOR_GREEN),
    ("Try at home", "Add Freeze/Unfreeze commands, transaction history in the view", COLOR_BLUE),
    ("Scaffold your own", "npx teob new aggregate — working code in seconds", COLOR_BLUE),
    ("Watch the recording", "Solution walkthroughs with chapter markers", COLOR_ORANGE),
    ("Next workshop", "Episode 2: Sagas & Cross-Aggregate Communication", COLOR_PURPLE),
]
for i, (title, desc, color) in enumerate(cta_items):
    y = Emu(2100000 + i * 800000)
    add_box(slide, M_LEFT, y, Emu(400000), Emu(400000), color, text="→", font_size=20)
    add_text(slide, Emu(1400000), y, Emu(5000000), Emu(400000),
             title, font_size=20, bold=True, color=COLOR_WHITE)
    add_text(slide, Emu(1400000), y + Emu(350000), Emu(10000000), Emu(350000),
             desc, font_size=15, color=COLOR_LIGHT_GRAY)

add_text(slide, M_LEFT, Emu(6200000), CONTENT_W, Emu(400000),
         "Thank you!  •  Questions?",
         font_size=24, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

add_notes(slide,
    "TIMING: 2:37–2:40\n"
    "\"You built an event-sourced service with a read model and AI integration.\"\n"
    "\"Three pure functions. That's the whole architecture.\"\n"
    "Open for questions."
)


# ══════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════

import os
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "DDD-Workshop-Slides.pptx")
prs.save(output_path)
print(f"Saved {len(prs.slides)} slides to {output_path}")
